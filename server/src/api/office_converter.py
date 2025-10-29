"""
Servizio di conversione documenti Office → PDF per preview
Mantiene il file originale intatto, crea solo una preview PDF
"""

import os
import subprocess
import tempfile
import shutil
from pathlib import Path
from typing import Optional, Tuple
import logging

logger = logging.getLogger(__name__)

class OfficeConverter:
    """Convertitore documenti Office → PDF per preview"""
    
    def __init__(self):
        self.temp_dir = tempfile.mkdtemp(prefix='securevox_office_')
        logger.info(f"OfficeConverter inizializzato con temp dir: {self.temp_dir}")
    
    def __del__(self):
        """Cleanup della directory temporanea"""
        try:
            if os.path.exists(self.temp_dir):
                shutil.rmtree(self.temp_dir)
                logger.info(f"Temp dir pulita: {self.temp_dir}")
        except Exception as e:
            logger.warning(f"Errore pulizia temp dir: {e}")
    
    def convert_to_pdf_preview(self, input_file_path: str, output_filename: str) -> Optional[str]:
        """
        Converte un documento Office in PDF per preview usando LibreOffice
        
        Args:
            input_file_path: Percorso del file Office originale
            output_filename: Nome del file PDF di output
            
        Returns:
            Percorso del PDF generato o None se errore
        """
        try:
            input_path = Path(input_file_path)
            if not input_path.exists():
                logger.error(f"File di input non trovato: {input_file_path}")
                return None
            
            # Usa LibreOffice per tutti i formati Office
            return self._convert_with_libreoffice(input_file_path, output_filename)
                
        except Exception as e:
            logger.error(f"Errore conversione {input_file_path}: {e}")
            return None
    
    def _convert_with_libreoffice(self, input_file_path: str, output_filename: str) -> Optional[str]:
        """
        Converte documento Office usando LibreOffice headless
        
        Args:
            input_file_path: Percorso del file Office da convertire
            output_filename: Nome del file PDF di output (senza estensione)
            
        Returns:
            Percorso del file PDF generato o None se fallisce
        """
        try:
            # Comando LibreOffice per conversione - prova diversi percorsi
            possible_paths = [
                '/Applications/LibreOffice.app/Contents/MacOS/soffice',  # macOS standard
                '/usr/bin/libreoffice',  # Linux standard
                '/usr/local/bin/libreoffice',  # Linux alternativo
                'libreoffice',  # PATH system
                'soffice',  # Alternativo
            ]
            
            libreoffice_path = None
            for path in possible_paths:
                if os.path.exists(path) or shutil.which(path):
                    libreoffice_path = path
                    break
            
            if not libreoffice_path:
                logger.error("❌ LibreOffice non trovato in nessun percorso standard")
                return self._create_fallback_pdf(output_filename, "Office Document")
            
            cmd = [
                libreoffice_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', self.temp_dir,
                input_file_path
            ]
            
            logger.info(f"🔄 Conversione Office→PDF con LibreOffice: {input_file_path}")
            
            # Esegui conversione
            result = subprocess.run(
                cmd,
                check=True,
                capture_output=True,
                text=True,
                timeout=60  # Timeout 60 secondi per LibreOffice
            )
            
            # LibreOffice crea il file con il nome originale + .pdf
            input_basename = os.path.splitext(os.path.basename(input_file_path))[0]
            generated_path = os.path.join(self.temp_dir, f"{input_basename}.pdf")
            output_path = os.path.join(self.temp_dir, f"{output_filename}.pdf")
            
            if os.path.exists(generated_path):
                # Rinomina al nome desiderato se diverso
                if generated_path != output_path:
                    shutil.move(generated_path, output_path)
                logger.info(f"✅ Conversione LibreOffice completata: {output_path}")
                return output_path
            else:
                logger.error(f"❌ File PDF non generato da LibreOffice: {generated_path}")
                return self._create_fallback_pdf(output_filename, "Office Document")
                
        except subprocess.TimeoutExpired:
            logger.error("❌ Timeout conversione LibreOffice")
            return self._create_fallback_pdf(output_filename, "Office Document")
        except subprocess.CalledProcessError as e:
            logger.error(f"❌ Errore conversione LibreOffice: {e}")
            logger.error(f"   stdout: {e.stdout}")
            logger.error(f"   stderr: {e.stderr}")
            return self._create_fallback_pdf(output_filename, "Office Document")
        except Exception as e:
            logger.error(f"❌ Errore inatteso conversione LibreOffice: {e}")
            return self._create_fallback_pdf(output_filename, "Office Document")
    
    def _convert_word_to_pdf(self, input_path: str, output_filename: str) -> Optional[str]:
        """Converte documento Word → PDF"""
        try:
            from docx import Document
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            
            logger.info(f"Conversione Word → PDF: {input_path}")
            
            # Leggi il documento Word
            doc = Document(input_path)
            
            # Crea PDF
            output_path = os.path.join(self.temp_dir, f"{output_filename}.pdf")
            pdf_doc = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            # Estrai il testo e convertilo in PDF
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    p = Paragraph(paragraph.text, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 12))
            
            if not story:
                # Se non c'è testo, crea una pagina con info file
                story.append(Paragraph(f"Documento Word: {output_filename}", styles['Title']))
                story.append(Spacer(1, 20))
                story.append(Paragraph("Preview del contenuto non disponibile", styles['Normal']))
            
            pdf_doc.build(story)
            logger.info(f"Word convertito con successo: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Errore conversione Word: {e}")
            return self._create_fallback_pdf(output_filename, "Word Document")
    
    def _convert_excel_to_pdf(self, input_path: str, output_filename: str) -> Optional[str]:
        """Converte foglio Excel → PDF"""
        try:
            from openpyxl import load_workbook
            from reportlab.lib.pagesizes import letter, landscape
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib import colors
            
            logger.info(f"Conversione Excel → PDF: {input_path}")
            
            # Leggi il foglio Excel
            wb = load_workbook(input_path, data_only=True)
            ws = wb.active
            
            # Crea PDF
            output_path = os.path.join(self.temp_dir, f"{output_filename}.pdf")
            pdf_doc = SimpleDocTemplate(output_path, pagesize=landscape(letter))
            styles = getSampleStyleSheet()
            story = []
            
            # Aggiungi titolo
            story.append(Paragraph(f"Foglio Excel: {output_filename}", styles['Title']))
            story.append(Paragraph("", styles['Normal']))  # Spacer
            
            # Estrai dati (max 10 righe e 10 colonne per preview)
            data = []
            max_rows = min(10, ws.max_row)
            max_cols = min(10, ws.max_column)
            
            for row in range(1, max_rows + 1):
                row_data = []
                for col in range(1, max_cols + 1):
                    cell_value = ws.cell(row=row, column=col).value
                    row_data.append(str(cell_value) if cell_value is not None else "")
                data.append(row_data)
            
            if data:
                # Crea tabella
                table = Table(data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                story.append(table)
            else:
                story.append(Paragraph("Nessun dato trovato nel foglio", styles['Normal']))
            
            pdf_doc.build(story)
            logger.info(f"Excel convertito con successo: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Errore conversione Excel: {e}")
            return self._create_fallback_pdf(output_filename, "Excel Spreadsheet")
    
    def _convert_powerpoint_to_pdf(self, input_path: str, output_filename: str) -> Optional[str]:
        """Converte presentazione PowerPoint → PDF"""
        try:
            from pptx import Presentation
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
            from reportlab.lib.styles import getSampleStyleSheet
            
            logger.info(f"Conversione PowerPoint → PDF: {input_path}")
            
            # Leggi la presentazione
            prs = Presentation(input_path)
            
            # Crea PDF
            output_path = os.path.join(self.temp_dir, f"{output_filename}.pdf")
            pdf_doc = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            # Titolo
            story.append(Paragraph(f"Presentazione PowerPoint: {output_filename}", styles['Title']))
            story.append(Spacer(1, 20))
            
            # Estrai testo dalle slide (max 5 slide per preview)
            slide_count = 0
            for slide in prs.slides:
                if slide_count >= 5:  # Limita a 5 slide per preview
                    break
                    
                slide_count += 1
                story.append(Paragraph(f"Slide {slide_count}", styles['Heading2']))
                story.append(Spacer(1, 12))
                
                # Estrai testo dalla slide
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    for text in slide_text:
                        story.append(Paragraph(text, styles['Normal']))
                        story.append(Spacer(1, 8))
                else:
                    story.append(Paragraph("(Slide senza testo)", styles['Italic']))
                
                if slide_count < len(prs.slides):
                    story.append(PageBreak())
            
            if slide_count == 0:
                story.append(Paragraph("Nessuna slide trovata", styles['Normal']))
            
            pdf_doc.build(story)
            logger.info(f"PowerPoint convertito con successo: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Errore conversione PowerPoint: {e}")
            return self._create_fallback_pdf(output_filename, "PowerPoint Presentation")
    
    def _create_fallback_pdf(self, filename: str, doc_type: str) -> str:
        """Crea un PDF di fallback quando la conversione fallisce"""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            
            output_path = os.path.join(self.temp_dir, f"{filename}_fallback.pdf")
            pdf_doc = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            
            story = [
                Paragraph(f"{doc_type}", styles['Title']),
                Spacer(1, 20),
                Paragraph(f"Nome file: {filename}", styles['Normal']),
                Spacer(1, 12),
                Paragraph("Preview del contenuto non disponibile", styles['Normal']),
                Spacer(1, 12),
                Paragraph("Il file originale è disponibile per il download", styles['Italic']),
            ]
            
            pdf_doc.build(story)
            logger.info(f"PDF fallback creato: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Errore creazione PDF fallback: {e}")
            return None

# Istanza globale del convertitore
office_converter = OfficeConverter()
